"""
Builds Sentinel_Deliverable.pptx — slide deck mapped to the project rubric.

Each content slide carries the rubric criterion + point value in the header,
so a grader can score directly off the deck. Numbers and file references come
from the deployed system, not vibes.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ── Brand palette ────────────────────────────────────────────────────────────
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT = RGBColor(0x2E, 0x86, 0xDE)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xFA)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MUTED = RGBColor(0x55, 0x5C, 0x6A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x27, 0xAE, 0x60)
AMBER = RGBColor(0xE6, 0x7E, 0x22)
RED = RGBColor(0xC0, 0x39, 0x2B)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG = RGBColor(0xCD, 0xD6, 0xF4)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

TOTAL_SLIDES = 13


# ── Drawing helpers ──────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, line=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_width:
            shape.line.width = line_width
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=12, color=TEXT_DARK,
                bullet_color=ACCENT, marker="■", body_color=TEXT_MUTED):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            head, body = item
        else:
            head, body = item, ""
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(5)
        r1 = p.add_run()
        r1.text = f"{marker}  "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = head
        r2.font.size = Pt(size)
        r2.font.bold = True
        r2.font.color.rgb = color
        if body:
            r3 = p.add_run()
            r3.text = "  " + body
            r3.font.size = Pt(size)
            r3.font.color.rgb = body_color
    return tb


def slide_header(slide, criterion, points, title, kicker=None):
    add_rect(slide, 0, 0, SW, Inches(0.55), NAVY)
    add_text(slide, Inches(0.4), 0, Inches(9), Inches(0.55),
             f"SENTINEL  •  {criterion}",
             size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(11.0), 0, Inches(2.0), Inches(0.55),
             f"Rubric: {points} pts", size=12, bold=True, color=ACCENT,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(0.55), Inches(0.7), Inches(12), Inches(0.7),
             title, size=28, bold=True, color=NAVY)
    add_rect(slide, Inches(0.55), Inches(1.4), Inches(0.7), Inches(0.06), ACCENT)
    if kicker:
        add_text(slide, Inches(0.55), Inches(1.5), Inches(12), Inches(0.4),
                 kicker, size=13, color=TEXT_MUTED)


def slide_footer(slide, page_num):
    add_rect(slide, 0, Inches(7.15), SW, Inches(0.35), LIGHT_BG)
    add_text(slide, Inches(0.4), Inches(7.15), Inches(8), Inches(0.35),
             "Sentinel  •  Risk Analysis Pipeline  •  v1.0.0",
             size=10, color=TEXT_MUTED, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.4), Inches(0.35),
             f"{page_num} / {TOTAL_SLIDES}",
             size=10, color=TEXT_MUTED, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)


def evidence_strip(slide, top_y, items):
    """Yellow-highlighted evidence strip at bottom of content area."""
    add_rect(slide, Inches(0.55), top_y, Inches(12.2), Inches(0.55), LIGHT_BG)
    add_rect(slide, Inches(0.55), top_y, Inches(0.12), Inches(0.55), GREEN)
    text = "  •  ".join(items)
    add_text(slide, Inches(0.85), top_y, Inches(11.8), Inches(0.55),
             "EVIDENCE  " + text, size=11, color=NAVY, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)


def code_block(slide, x, y, w, h, lines, *, size=10):
    """Monospace code-style block."""
    add_rect(slide, x, y, w, h, CODE_BG)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.name = "Consolas"
        r.font.size = Pt(size)
        r.font.color.rgb = CODE_FG


# ============================================================================
# 1. TITLE
# ============================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(3.55), Inches(1.6), Inches(0.12), ACCENT)
add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(0.6),
         "SENTINEL", size=20, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(2.85), Inches(12), Inches(1.2),
         "Risk Analysis Pipeline", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.9), Inches(12), Inches(0.6),
         "AI-assisted, deterministic-decision underwriting for loan documents.",
         size=18, color=RGBColor(0xCF, 0xD8, 0xE3))
add_text(s, Inches(0.7), Inches(4.3), Inches(12), Inches(0.6),
         "Deployed on Google Cloud Run · v1.0.0",
         size=18, color=RGBColor(0xCF, 0xD8, 0xE3))
add_text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.4),
         "Project Deliverable  •  CISC 867010 — University of Delaware",
         size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.4),
         "Industry partner: Best Egg  ·  github.com/Khey17/sentinel-nextgenai-execution-layer",
         size=12, color=RGBColor(0x9F, 0xAD, 0xC0))


# ============================================================================
# 2. PROJECT OVERVIEW, GOALS, SCOPE  (10 pts)
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Project Overview, Goals, Scope", "10",
             "What Best Egg Asked For — and What We Built",
             kicker="The brief, the scope boundary, the stakeholders, the guarantees")

add_rect(s, Inches(0.55), Inches(2.0), Inches(12.2), Inches(1.0), LIGHT_BG)
add_rect(s, Inches(0.55), Inches(2.0), Inches(0.12), Inches(1.0), ACCENT)
add_text(s, Inches(0.85), Inches(2.1), Inches(11.8), Inches(0.85),
         "“Help define and create a scalable, resilient and secure GenAI "
         "infrastructure and orchestration layer to integrate more "
         "unstructured data… in a scalable, resilient, and secure way.”",
         size=14, color=NAVY)

# Two columns: In-scope, Out-of-scope
add_rect(s, Inches(0.55), Inches(3.2), Inches(6.0), Inches(3.5), WHITE,
         line=RGBColor(0xE3, 0xE8, 0xEF))
add_rect(s, Inches(0.55), Inches(3.2), Inches(6.0), Inches(0.5), GREEN)
add_text(s, Inches(0.85), Inches(3.2), Inches(5.4), Inches(0.5),
         "In scope (MVP)", size=14, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(0.85), Inches(3.85), Inches(5.6), Inches(2.8), [
    ("Digital PDF documents", "bank statements, paystubs, W-2s, tax returns"),
    ("End-to-end pipeline", "upload → parse → redact → extract → score → audit"),
    ("Reason-coded decisions", "every PASS / NEEDS_REVIEW with named codes"),
    ("Application as a unit", "merge per-doc evidence into one decision"),
    ("Cloud deployment", "GCP Cloud Run, GitHub Actions auto-deploy"),
], size=11, marker="✓", bullet_color=GREEN)

add_rect(s, Inches(6.75), Inches(3.2), Inches(6.0), Inches(3.5), WHITE,
         line=RGBColor(0xE3, 0xE8, 0xEF))
add_rect(s, Inches(6.75), Inches(3.2), Inches(6.0), Inches(0.5), AMBER)
add_text(s, Inches(7.05), Inches(3.2), Inches(5.4), Inches(0.5),
         "Out of scope (deliberately)", size=14, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(7.05), Inches(3.85), Inches(5.6), Inches(2.8), [
    ("OCR for scanned PDFs", "digital text only — image PDFs deferred"),
    ("Email / chat / image ingestion", "PDF upload only"),
    ("Lending decision automation", "we recommend; underwriters decide"),
    ("Multi-tenant auth", "single-tenant MVP, IAM as a follow-on"),
    ("Phase 3 agentic pipeline", "scrapped — deterministic > orchestration"),
], size=11, marker="—", bullet_color=AMBER)

evidence_strip(s, Inches(6.85), [
    "README.md (200+ lines)",
    "Phase 0/1/2 status table",
    "Stakeholder matrix",
    "Privacy & auditability guarantees",
])
slide_footer(s, 2)


# ============================================================================
# 3. TECHNICAL APPROACH & SYSTEM DESIGN  (15 pts)
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Technical Approach & System Design", "15",
             "Deployed Cloud Topology + Pipeline",
             kicker="Three Cloud Run services + Memorystore + Cloud SQL + GCS + Vertex AI")

# Top row: deployed services
y_top = Inches(2.05)
def node(x, y, w, h, label, *, fill=WHITE, text_color=NAVY, border=ACCENT, size=11):
    sh = add_rect(s, x, y, w, h, fill, line=border, line_width=Pt(1.25))
    add_text(s, x, y, w, h, label, size=size, bold=True, color=text_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return sh

def arrow(x1, y1, x2, y2):
    line = s.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = TEXT_MUTED
    line.line.width = Pt(1.25)
    ln = line.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")

node(Inches(0.55), y_top, Inches(2.0), Inches(0.65),
     "Streamlit\nCloud Run", fill=NAVY, text_color=WHITE, border=NAVY)
node(Inches(3.05), y_top, Inches(2.2), Inches(0.65),
     "FastAPI\nCloud Run · 2Gi", fill=ACCENT, text_color=WHITE, border=ACCENT)
node(Inches(5.75), y_top, Inches(2.0), Inches(0.65),
     "Cloud SQL\nPostgres")
node(Inches(8.25), y_top, Inches(2.0), Inches(0.65),
     "Cloud Storage\nbucket")
node(Inches(10.75), y_top, Inches(2.05), Inches(0.65),
     "Memorystore\nRedis")

mid = y_top + Inches(0.325)
arrow(Inches(2.55), mid, Inches(3.05), mid)
arrow(Inches(5.25), mid, Inches(5.75), mid)
arrow(Inches(7.75), mid, Inches(8.25), mid)
arrow(Inches(10.25), mid, Inches(10.75), mid)

# Worker
node(Inches(10.75), Inches(3.15), Inches(2.05), Inches(0.65),
     "Celery Worker\nmin-instances=1", fill=NAVY, text_color=WHITE,
     border=NAVY, size=10)
arrow(Inches(11.78), y_top + Inches(0.65), Inches(11.78), Inches(3.15))

# Pipeline strip
add_text(s, Inches(0.55), Inches(4.05), Inches(8), Inches(0.4),
         "Per-document pipeline (sequential, idempotent steps)",
         size=12, bold=True, color=NAVY)
add_rect(s, Inches(0.55), Inches(4.45), Inches(12.2), Inches(1.05), LIGHT_BG)
steps = [
    ("Parse", "pdfplumber"),
    ("Classify", "Keyword gate"),
    ("Authenticate", "Math + meta"),
    ("Redact", "Presidio + spaCy"),
    ("Extract", "Gemini 2.5 Flash"),
    ("Score", "Deterministic"),
]
step_w = Inches(1.85); gap = Inches(0.12)
start_x = Inches(0.7)
for i, (h, sub) in enumerate(steps):
    x = start_x + i * (step_w + gap)
    node(x, Inches(4.6), step_w, Inches(0.85), f"{h}\n{sub}", size=10)
    if i < len(steps) - 1:
        arrow(x + step_w, Inches(5.025), x + step_w + gap, Inches(5.025))

# Application merge
add_rect(s, Inches(0.55), Inches(5.7), Inches(12.2), Inches(0.95), WHITE,
         line=GREEN, line_width=Pt(1.5))
add_rect(s, Inches(0.55), Inches(5.7), Inches(0.12), Inches(0.95), GREEN)
add_text(s, Inches(0.85), Inches(5.75), Inches(11.8), Inches(0.4),
         "Application-level merge  ·  GET /batches/{id}/decision",
         size=13, bold=True, color=GREEN)
add_text(s, Inches(0.85), Inches(6.15), Inches(11.8), Inches(0.5),
         "Merge per-doc extractions: income from paystub > W-2 > tax return > "
         "bank deposits; risk signals from bank statement; auth = min across docs.",
         size=10, color=TEXT_MUTED)

evidence_strip(s, Inches(6.7), [
    ".github/workflows/deploy.yml",
    "src/api/app/aggregator.py",
    "src/api/app/worker.py",
    "ARCHITECTURE.md",
])
slide_footer(s, 3)


# ============================================================================
# 4. IMPLEMENTATION QUALITY & CODE ORGANIZATION  (10 pts)
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Implementation Quality & Code Organization", "10",
             "One Module, One Job — Clear Separation of Concerns",
             kicker="Repository layout + naming + type hints + minimal duplication")

# Left: repo tree (code style)
code_block(s, Inches(0.55), Inches(2.05), Inches(6.4), Inches(4.5), [
    "sentinel-nextgenai-execution-layer/",
    "├── src/api/",
    "│   ├── app/",
    "│   │   ├── main.py          # FastAPI app",
    "│   │   ├── routes.py        # HTTP endpoints",
    "│   │   ├── worker.py        # Celery tasks",
    "│   │   ├── models.py        # SQLAlchemy ORM",
    "│   │   ├── schemas.py       # Pydantic I/O",
    "│   │   ├── extractor.py     # Vertex AI Gemini",
    "│   │   ├── redactor.py      # Presidio + spaCy",
    "│   │   ├── authenticator.py # PDF math + meta",
    "│   │   ├── scorer.py        # Deterministic 100pt",
    "│   │   ├── aggregator.py    # Application merge",
    "│   │   ├── guardrails.py    # Edge validation",
    "│   │   ├── storage.py       # GCS / MinIO",
    "│   │   └── metrics.py       # Prometheus",
    "│   ├── Dockerfile",
    "│   └── requirements.txt",
    "├── frontend/app.py          # Streamlit",
    "├── tests/                   # 4 unit modules",
    "├── monitoring/              # Prometheus + Grafana",
    "└── .github/workflows/       # CI + Deploy",
], size=11)

# Right: implementation principles
add_rect(s, Inches(7.15), Inches(2.05), Inches(5.6), Inches(4.5), WHITE,
         line=RGBColor(0xE3, 0xE8, 0xEF))
add_rect(s, Inches(7.15), Inches(2.05), Inches(5.6), Inches(0.5), ACCENT)
add_text(s, Inches(7.45), Inches(2.05), Inches(5.0), Inches(0.5),
         "Conventions enforced", size=14, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)

add_bullets(s, Inches(7.45), Inches(2.7), Inches(5.2), Inches(3.8), [
    ("One module, one responsibility",
     "scorer.py never calls Gemini; redactor.py never touches the DB."),
    ("Type hints throughout",
     "from __future__ import annotations in every module; PEP 604 unions."),
    ("Pydantic at I/O boundaries",
     "schemas.py contains every request/response shape — no untyped dicts in routes."),
    ("Idempotent worker tasks",
     "Each Celery task overwrites its output artifact — safe to re-run."),
    ("Comments explain WHY, not what",
     "scorer.py docstring: “Not 'AI said 0.74' — but reason codes that ship.”"),
    ("Minimal duplication",
     "compute_score() reused for per-doc and application; merge_extractions does not duplicate scoring logic."),
], size=11)

evidence_strip(s, Inches(6.7), [
    "16 modules in src/api/app",
    "100% type-annotated public APIs",
    "Pydantic schemas at all routes",
    "scorer.py reused via aggregator",
])
slide_footer(s, 4)


# ============================================================================
# 5. ML MODELS, METRICS & PROFILING  (10 pts — Performance)
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "ML Models, Metrics & Profiling", "10",
             "What the System Runs and How We Measure It",
             kicker="Three models in production · Prometheus instrumented · documented choice of metric")

# Top: ML models table
top = Inches(2.0)
add_rect(s, Inches(0.55), top, Inches(12.2), Inches(0.5), NAVY)
add_text(s, Inches(0.85), top, Inches(3.5), Inches(0.5),
         "Model", size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(4.45), top, Inches(2.2), Inches(0.5),
         "Role", size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(6.75), top, Inches(2.0), Inches(0.5),
         "Latency (p50)", size=12, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(8.85), top, Inches(3.85), Inches(0.5),
         "Why this model", size=12, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)

ml_rows = [
    ("Gemini 2.5 Flash",
     "Structured field extraction (Vertex AI, schema-pinned)",
     "5 – 20 s",
     "Fast tier, structured-output JSON schema, region-local"),
    ("spaCy en_core_web_lg",
     "Named-entity recognition for PII redaction",
     "200 – 500 ms",
     "Higher-recall NER than _sm; runs CPU-only in 2 GiB"),
    ("Presidio analyzer",
     "Pattern-based PII (SSN, accounts, routing #)",
     "100 – 300 ms",
     "Deterministic regex tier — catches what spaCy misses"),
]
for i, (model, role, lat, why) in enumerate(ml_rows):
    y = top + Inches(0.5) + i * Inches(0.5)
    bg = WHITE if i % 2 == 0 else LIGHT_BG
    add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.5), bg)
    add_text(s, Inches(0.85), y, Inches(3.5), Inches(0.5),
             model, size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.45), y, Inches(2.2), Inches(0.5),
             role, size=10, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(6.75), y, Inches(2.0), Inches(0.5),
             lat, size=11, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.85), y, Inches(3.85), Inches(0.5),
             why, size=10, color=TEXT_MUTED, anchor=MSO_ANCHOR.MIDDLE)

# Bottom: profiling + metric choice
y_bot = Inches(4.05)
add_rect(s, Inches(0.55), y_bot, Inches(6.0), Inches(2.6), WHITE,
         line=RGBColor(0xE3, 0xE8, 0xEF))
add_rect(s, Inches(0.55), y_bot, Inches(6.0), Inches(0.5), ACCENT)
add_text(s, Inches(0.85), y_bot, Inches(5.4), Inches(0.5),
         "Profiling — how we know", size=13, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(0.85), y_bot + Inches(0.6), Inches(5.6), Inches(2.0), [
    ("Per-step timers", "worker.py records start = time.time(), record_step(name, elapsed)"),
    ("Prometheus metrics", "histograms per step via prometheus-fastapi-instrumentator"),
    ("Grafana 15-panel board", "throughput, per-step latency, failure-by-step, queue depth"),
    ("Synthetic load harness", "tests/generate_1000_pairs.py for stress testing"),
], size=10)

add_rect(s, Inches(6.75), y_bot, Inches(6.0), Inches(2.6), WHITE,
         line=RGBColor(0xE3, 0xE8, 0xEF))
add_rect(s, Inches(6.75), y_bot, Inches(6.0), Inches(0.5), ACCENT)
add_text(s, Inches(7.05), y_bot, Inches(5.4), Inches(0.5),
         "Metric choice — and why", size=13, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(7.05), y_bot + Inches(0.6), Inches(5.6), Inches(2.0), [
    ("Deterministic 100-pt score", "Replaces 'LLM said 0.74' with reason codes"),
    ("Threshold = 0.80", "Configurable via CONFIDENCE_THRESHOLD env var"),
    ("Hard stops bypass score", "AUTH_FAILED / INTEGRITY_FAIL force review at any score"),
    ("Why not LLM confidence?", "Not auditable, drifts across model versions, no reasons"),
], size=10)

evidence_strip(s, Inches(6.7), [
    "scorer.py · 100pt scorecard with reason codes",
    "monitoring/grafana/ · 15-panel board",
    "metrics.py · Prometheus instruments",
])
slide_footer(s, 5)


# ============================================================================
# 6. PERFORMANCE, RESOURCES & SCALABILITY  (10 pts)
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Performance, Resources & Scalability", "10",
             "Honest Numbers from the Deployed System",
             kicker="Per-step latency, steady-state throughput, hard ceilings")

# Per-step latency table
top = Inches(2.05)
add_rect(s, Inches(0.55), top, Inches(7.2), Inches(0.5), NAVY)
add_text(s, Inches(0.85), top, Inches(2.8), Inches(0.5),
         "Step", size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(3.7), top, Inches(1.5), Inches(0.5),
         "Latency", size=12, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(5.2), top, Inches(2.5), Inches(0.5),
         "Bottleneck", size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

step_rows = [
    ("Parse (pdfplumber)",        "1 – 3 s",   "PDF page count"),
    ("Classify (keyword)",        "< 1 s",     "Pure CPU"),
    ("Authenticate (math+meta)",  "1 – 2 s",   "Deterministic"),
    ("Redact (Presidio+spaCy)",   "2 – 5 s",   "spaCy en_core_web_lg"),
    ("Extract (Gemini Flash)",    "5 – 20 s",  "Vertex AI round-trip"),
    ("Score + cleanup",           "< 1 s",     "Pure Python"),
]
for i, (step, lat, bn) in enumerate(step_rows):
    y = top + Inches(0.5) + i * Inches(0.4)
    bg = WHITE if i % 2 == 0 else LIGHT_BG
    add_rect(s, Inches(0.55), y, Inches(7.2), Inches(0.4), bg)
    add_text(s, Inches(0.85), y, Inches(2.8), Inches(0.4),
             step, size=11, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.7), y, Inches(1.5), Inches(0.4),
             lat, size=11, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.2), y, Inches(2.5), Inches(0.4),
             bn, size=10, color=TEXT_MUTED, anchor=MSO_ANCHOR.MIDDLE)

# Total
y_total = top + Inches(0.5) + 6 * Inches(0.4) + Inches(0.05)
add_rect(s, Inches(0.55), y_total, Inches(7.2), Inches(0.5), NAVY)
add_text(s, Inches(0.85), y_total, Inches(2.8), Inches(0.5),
         "Total per document", size=12, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(3.7), y_total, Inches(1.5), Inches(0.5),
         "10 – 30 s", size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(5.2), y_total, Inches(2.5), Inches(0.5),
         "Gemini dominates", size=11, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# Right: scalability scenarios
add_rect(s, Inches(8.05), Inches(2.05), Inches(4.7), Inches(4.55), WHITE,
         line=RGBColor(0xE3, 0xE8, 0xEF))
add_rect(s, Inches(8.05), Inches(2.05), Inches(4.7), Inches(0.5), ACCENT)
add_text(s, Inches(8.05), Inches(2.05), Inches(4.7), Inches(0.5),
         "1,000-user load — honest answer", size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

scenarios = [
    ("/ 1 day", "1.4 docs/min", "✓ today", GREEN),
    ("/ 1 hour", "33 docs/min", "△ scale-out", AMBER),
    ("/ 5 min", "400 docs/min", "✗ Gemini cap", RED),
]
y_sc = Inches(2.7)
for i, (when, demand, verdict, color) in enumerate(scenarios):
    yy = y_sc + i * Inches(0.65)
    add_rect(s, Inches(8.25), yy, Inches(4.3), Inches(0.55), LIGHT_BG)
    add_text(s, Inches(8.45), yy, Inches(1.5), Inches(0.55),
             when, size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.95), yy, Inches(1.6), Inches(0.55),
             demand, size=11, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(11.55), yy, Inches(0.95), Inches(0.55),
             verdict, size=11, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(8.25), Inches(4.85), Inches(4.3), Inches(1.6),
         "Hard ceiling: Vertex AI Gemini 2.5 Flash quota ≈ 300 RPM/region. "
         "Per-worker rate_limit=\"100/m\" is a safety throttle, not the ceiling. "
         "Even with infinite workers, Gemini caps the system.",
         size=10, color=TEXT_MUTED)

evidence_strip(s, Inches(6.7), [
    "5 docs/min/worker · 300/hr · 7,200/day",
    "Worker: 2 CPU / 2 GiB · concurrency=1",
    "Documented in deploy.yml + entrypoint.sh",
])
slide_footer(s, 6)


# ============================================================================
# 7. TESTING, VALIDATION & CORRECTNESS  (15 pts)
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Testing, Validation & Correctness", "15",
             "How We Know the Software Works",
             kicker="688 lines of unit tests · CI on every push · documented edge cases")

# Two columns
add_rect(s, Inches(0.55), Inches(2.05), Inches(6.0), Inches(4.95), WHITE,
         line=RGBColor(0xE3, 0xE8, 0xEF))
add_rect(s, Inches(0.55), Inches(2.05), Inches(6.0), Inches(0.5), NAVY)
add_text(s, Inches(0.85), Inches(2.05), Inches(5.4), Inches(0.5),
         "Unit tests — core logic", size=14, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(0.85), Inches(2.7), Inches(5.6), Inches(2.5), [
    ("test_authenticator.py · 129 lines",
     "PDF balance reconciliation, document type classification, fraud flags"),
    ("test_guardrails.py · 133 lines",
     "MIME, magic bytes, file size, PII-dump rejection at the edge"),
    ("test_classifier.py · 213 lines",
     "Keyword gate: financial vs non-financial documents, edge cases"),
    ("test_redactor.py · 213 lines",
     "Presidio + spaCy: SSN, accounts, names; placeholder consistency"),
], size=10)

add_text(s, Inches(0.85), Inches(5.3), Inches(5.4), Inches(0.4),
         "Edge cases & failure modes", size=12, bold=True, color=ACCENT)
add_bullets(s, Inches(0.85), Inches(5.7), Inches(5.6), Inches(1.4), [
    ("Empty PDF, corrupt PDF", "guardrails reject before pipeline starts"),
    ("Non-financial document", "classifier rejects (lease, receipt → FAILED)"),
    ("PII leakage in LLM output", "post-extraction PII scan + retry block"),
    ("Duplicate file upload", "SHA-256 dedup at /batches/upload"),
], size=10, marker="⚠")

# Right column: integration + CI
add_rect(s, Inches(6.75), Inches(2.05), Inches(6.0), Inches(4.95), WHITE,
         line=RGBColor(0xE3, 0xE8, 0xEF))
add_rect(s, Inches(6.75), Inches(2.05), Inches(6.0), Inches(0.5), NAVY)
add_text(s, Inches(7.05), Inches(2.05), Inches(5.4), Inches(0.5),
         "Integration / E2E + CI", size=14, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)

add_bullets(s, Inches(7.05), Inches(2.7), Inches(5.6), Inches(2.0), [
    ("Live deployment as E2E",
     "sentinel-frontend-1041799394320.us-central1.run.app — every push tests the full pipeline"),
    ("Application merge verified",
     "paystub + bank statement → 100/100 PASS, locally reproduced"),
    ("Smoke flow",
     "upload → status poll → results fetch — exercised on every demo"),
    ("Bulk fixtures",
     "tests/generate_1000_pairs.py · synthesizes paired documents at scale"),
], size=10)

# CI evidence box
add_rect(s, Inches(7.05), Inches(4.85), Inches(5.4), Inches(2.0), CODE_BG)
add_text(s, Inches(7.25), Inches(4.95), Inches(5.0), Inches(0.4),
         "CI runner — .github/workflows/ci.yml",
         size=11, bold=True, color=CODE_FG)
code_block(s, Inches(7.05), Inches(5.35), Inches(5.4), Inches(1.5), [
    "- run: pip install -r src/api/requirements.txt",
    "- run: python -m spacy download en_core_web_lg",
    "- run: pytest tests/ -v --tb=short",
], size=9)

evidence_strip(s, Inches(6.85), [
    "tests/ · 4 modules · 688 LOC",
    ".github/workflows/ci.yml on every push",
    "Live URL exercises full pipeline",
])
slide_footer(s, 7)


# ============================================================================
# 8. REPRODUCIBILITY & EXECUTION  (15 pts)
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Reproducibility & Execution", "15",
             "Two Targets, Pinned Environment, One Quickstart",
             kicker="Local Docker Compose for development · GCP Cloud Run for production")

# Local quickstart (code block)
add_text(s, Inches(0.55), Inches(2.05), Inches(6), Inches(0.4),
         "Local — Docker Compose", size=14, bold=True, color=NAVY)
code_block(s, Inches(0.55), Inches(2.5), Inches(6.4), Inches(2.6), [
    "# Clone and start the stack (~3 min on first run)",
    "git clone https://github.com/Khey17/sentinel-\\",
    "    nextgenai-execution-layer.git",
    "cd sentinel-nextgenai-execution-layer",
    "cp .env.example .env",
    "docker compose up --build",
    "",
    "# Frontend:  http://localhost:8501",
    "# API docs:  http://localhost:8000/docs",
    "# Grafana:   http://localhost:3000",
], size=10)

# GCP deploy (code block)
add_text(s, Inches(7.0), Inches(2.05), Inches(6), Inches(0.4),
         "Production — GCP Cloud Run", size=14, bold=True, color=NAVY)
code_block(s, Inches(7.0), Inches(2.5), Inches(5.8), Inches(2.6), [
    "# Triggered by push to main",
    "git push origin main",
    "",
    "# .github/workflows/deploy.yml builds + deploys:",
    "#   sentinel-api      (Cloud Run, 2Gi)",
    "#   sentinel-worker   (Cloud Run, min-instances=1)",
    "#   sentinel-frontend (Cloud Run)",
    "",
    "# Live: sentinel-frontend-...us-central1.run.app",
], size=10)

# Bottom strip: pinning + I/O
add_rect(s, Inches(0.55), Inches(5.25), Inches(12.25), Inches(1.4), WHITE,
         line=RGBColor(0xE3, 0xE8, 0xEF))
add_rect(s, Inches(0.55), Inches(5.25), Inches(0.12), Inches(1.4), GREEN)
add_text(s, Inches(0.85), Inches(5.35), Inches(11.8), Inches(0.4),
         "Environment is pinned and deterministic", size=13, bold=True, color=NAVY)
add_bullets(s, Inches(0.85), Inches(5.75), Inches(11.8), Inches(0.85), [
    ("requirements.txt", "exact versions for fastapi, sqlalchemy, psycopg, celery, presidio, spacy"),
    ("Dockerfile", "python:3.11-slim base, model downloaded at build time, deterministic across hosts"),
    ("docker-compose.yml", "MinIO, Postgres, Redis, API, Worker, Frontend, Prometheus, Grafana"),
], size=10)

evidence_strip(s, Inches(6.7), [
    "README.md quickstart",
    "src/api/requirements.txt (pinned)",
    ".env.example",
    "docker-compose.yml (8 services)",
    "Verified on macOS + Cloud Run Linux",
])
slide_footer(s, 8)


# ============================================================================
# 9. DOCUMENTATION, LICENSING & RELEASE READINESS  (10 pts)
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Documentation, Licensing & Release Readiness", "10",
             "Files a Grader Can Open",
             kicker="README · LICENSE · CHANGELOG · RELEASE_NOTES · pyproject · v1.0.0 tag")

# Top half: file matrix
top = Inches(2.05)
files = [
    ("README.md", "Project goals, scope, stakeholders, quickstart, architecture diagram"),
    ("LICENSE", "Apache License 2.0 (with NSF-DARSE attribution)"),
    ("NOTICE", "Third-party attributions"),
    ("CHANGELOG.md", "Keep-a-Changelog format · v1.0.0 entry with all features and fixes"),
    ("RELEASE_NOTES.md", "v1.0.0 high-level summary, known issues, migration notes"),
    ("pyproject.toml", "Package metadata, version, pinned deps, optional [test] extras"),
    ("ARCHITECTURE.md", "Mermaid diagram + component-by-component description"),
    ("docs/", "Sphinx scaffold + course PDF artifact"),
]
add_rect(s, Inches(0.55), top, Inches(12.2), Inches(0.5), NAVY)
add_text(s, Inches(0.85), top, Inches(3.5), Inches(0.5),
         "Artifact", size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(4.65), top, Inches(8), Inches(0.5),
         "Purpose", size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
for i, (name, desc) in enumerate(files):
    y = top + Inches(0.5) + i * Inches(0.42)
    bg = WHITE if i % 2 == 0 else LIGHT_BG
    add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.42), bg)
    add_text(s, Inches(0.85), y, Inches(3.5), Inches(0.42),
             name, size=11, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.65), y, Inches(8), Inches(0.42),
             desc, size=10, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)

# Versioning callout
y_v = Inches(6.0)
add_rect(s, Inches(0.55), y_v, Inches(12.2), Inches(0.7), WHITE,
         line=GREEN, line_width=Pt(1.5))
add_rect(s, Inches(0.55), y_v, Inches(0.12), Inches(0.7), GREEN)
add_text(s, Inches(0.85), y_v, Inches(11.8), Inches(0.7),
         "Versioning: SemVer · current = v1.0.0 (tagged release) · "
         "auto-deploy from main on every push",
         size=13, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

evidence_strip(s, Inches(6.85), [
    "Apache 2.0 LICENSE",
    "v1.0.0 git tag + GH release",
    "CHANGELOG follows Keep-a-Changelog",
    "pyproject.toml installable",
])
slide_footer(s, 9)


# ============================================================================
# 10. RESULTS, DELIVERABLES & PROJECT-SPECIFIC OUTCOMES  (10 pts)
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Results, Deliverables & Outcomes", "10",
             "What Best Egg Asked For — and What Now Exists",
             kicker="Each line item: deliverable from the brief vs the artifact that ships it")

deliverables = [
    ("Framework",
     "End-to-end pipeline upload → parse → redact → extract → score → audit",
     "src/api/app/* · 16 modules · 5 worker tasks · idempotent"),
    ("Visualization Dashboard",
     "Customer + Business views with reason codes, score breakdown, redaction proof",
     "frontend/app.py · live at sentinel-frontend...run.app"),
    ("Infrastructure",
     "Running, auto-deploying system on GCP — not a prototype sketch",
     ".github/workflows/deploy.yml · 3 Cloud Run services · pinned worker"),
    ("Privacy guarantee",
     "PII never reaches the LLM (enforced by pipeline order, not discipline)",
     "redactor.py runs before extractor.py in chain · Presidio + spaCy"),
    ("Reason-coded explainability",
     "Every PASS / NEEDS_REVIEW carries the specific reason codes that justified it",
     "scorer.py · ECOA / Reg B-ready adverse action codes"),
    ("Application-level decision",
     "Bank statement + paystub merged into ONE recommendation per applicant",
     "aggregator.py · GET /batches/{id}/decision (NEW)"),
]

top = Inches(2.0)
row_h = Inches(0.78)
for i, (deliv, what, evidence) in enumerate(deliverables):
    y = top + i * (row_h + Inches(0.05))
    add_rect(s, Inches(0.55), y, Inches(12.2), row_h, WHITE,
             line=RGBColor(0xE3, 0xE8, 0xEF))
    add_rect(s, Inches(0.55), y, Inches(0.12), row_h, GREEN)
    add_rect(s, Inches(0.85), y + Inches(0.13), Inches(2.6), Inches(0.5), NAVY)
    add_text(s, Inches(0.85), y + Inches(0.13), Inches(2.6), Inches(0.5),
             deliv, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.65), y + Inches(0.08), Inches(9.0), Inches(0.4),
             what, size=11, color=TEXT_DARK)
    add_text(s, Inches(3.65), y + Inches(0.42), Inches(9.0), Inches(0.36),
             "→ " + evidence, size=9, color=TEXT_MUTED)

slide_footer(s, 10)


# ============================================================================
# 11. RELEASE NOTES  (5 pts — Documentation for releases)
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Release Notes — v1.0.0", "10",
             "What's in This Release",
             kicker="Initial production release · auto-deploy from main on push")

# Three columns: Features / Fixes / Known Issues
col_w = Inches(3.95)
top = Inches(2.05)

# Features
add_rect(s, Inches(0.55), top, col_w, Inches(4.65), WHITE,
         line=RGBColor(0xE3, 0xE8, 0xEF))
add_rect(s, Inches(0.55), top, col_w, Inches(0.5), GREEN)
add_text(s, Inches(0.55), top, col_w, Inches(0.5),
         "✦  Features", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(0.75), top + Inches(0.65), col_w - Inches(0.4), Inches(3.9), [
    ("Application-level merge",
     "/batches/{id}/decision combines paystub + bank statement into one decision"),
    ("Reason-coded customer view",
     "Score, risk band, top findings, ECOA / Reg B explainability"),
    ("Business review queue",
     "Score breakdown per item, one-click approve/reject with auto-attached reasons"),
    ("PII redaction",
     "Presidio + spaCy en_core_web_lg before any LLM call"),
    ("Auto-deploy",
     "GitHub Actions → Cloud Run on push to main"),
], size=10, marker="●")

# Fixes
add_rect(s, Inches(4.7), top, col_w, Inches(4.65), WHITE,
         line=RGBColor(0xE3, 0xE8, 0xEF))
add_rect(s, Inches(4.7), top, col_w, Inches(0.5), ACCENT)
add_text(s, Inches(4.7), top, col_w, Inches(0.5),
         "🛠  Fixes", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(4.9), top + Inches(0.65), col_w - Inches(0.4), Inches(3.9), [
    ("Worker scale-to-zero",
     "Cloud Run was reaping the worker (commit 92288a7); pinned --min-instances=1"),
    ("Reason codes surfaced",
     "Frontend was hiding the scorecard (commit 5324095); now rendered with severity"),
    ("Worker stability",
     "concurrency=1, prefetch=1 prevents head-of-line blocking (commit e2003c7)"),
    ("Null-income guard",
     "TypeError on $0/null monthly_net (commit 15d3639)"),
    ("Frontend poll cap",
     "2-min cap prevents infinite spinner on a stuck pipeline"),
], size=10, marker="●", bullet_color=ACCENT)

# Known issues
add_rect(s, Inches(8.85), top, col_w, Inches(4.65), WHITE,
         line=RGBColor(0xE3, 0xE8, 0xEF))
add_rect(s, Inches(8.85), top, col_w, Inches(0.5), AMBER)
add_text(s, Inches(8.85), top, col_w, Inches(0.5),
         "⚠  Known Issues", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(9.05), top + Inches(0.65), col_w - Inches(0.4), Inches(3.9), [
    ("Vertex AI quota cap",
     "300 RPM/region/project limits burst capacity"),
    ("Single worker bottleneck",
     "concurrency=1 trades throughput for memory headroom"),
    ("OCR not supported",
     "Scanned/image PDFs deferred — digital text only"),
    ("No multi-tenant auth",
     "Single-tenant MVP; IAM is a follow-on"),
    ("No formal migration",
     "v1.0.0 is initial release — no upgrade path needed yet"),
], size=10, marker="●", bullet_color=AMBER)

evidence_strip(s, Inches(6.85), [
    "CHANGELOG.md (Keep-a-Changelog format)",
    "RELEASE_NOTES.md (this slide)",
    "git tag v1.0.0 + GitHub Release",
])
slide_footer(s, 11)


# ============================================================================
# 12. REFLECTION & FUTURE WORK  (5 pts)
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Reflection & Future Work", "5",
             "What We Learned · What We'd Build Next",
             kicker="Honest retrospective — choices, regrets, the next milestone")

# Left: what worked / what didn't
add_rect(s, Inches(0.55), Inches(2.05), Inches(6.0), Inches(4.95), WHITE,
         line=RGBColor(0xE3, 0xE8, 0xEF))
add_rect(s, Inches(0.55), Inches(2.05), Inches(6.0), Inches(0.5), NAVY)
add_text(s, Inches(0.85), Inches(2.05), Inches(5.4), Inches(0.5),
         "Reflections", size=14, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.85), Inches(2.65), Inches(5.5), Inches(0.4),
         "What worked", size=12, bold=True, color=GREEN)
add_bullets(s, Inches(0.85), Inches(3.05), Inches(5.6), Inches(1.6), [
    ("Deterministic scorecard",
     "Replacing 'AI confidence' with reason codes was the single highest-leverage decision"),
    ("Scrapping Phase 3 agentic pipeline",
     "Determinism beat orchestration complexity; recovered 2 weeks"),
    ("Application-level merge",
     "Per-document scoring gave conflicting verdicts; merge was a clean fix"),
], size=10, marker="✓", bullet_color=GREEN)

add_text(s, Inches(0.85), Inches(4.85), Inches(5.5), Inches(0.4),
         "What didn't", size=12, bold=True, color=AMBER)
add_bullets(s, Inches(0.85), Inches(5.25), Inches(5.6), Inches(1.6), [
    ("Cloud Run for long-running workers",
     "Scale-to-zero burned us; --min-instances=1 fixed it but the pattern is wrong long-term"),
    ("Frontend hid the scorecard initially",
     "Hardcoded 'Low Risk' label undermined the deterministic backend"),
    ("CI was a no-op for too long",
     "ls -R . is not testing"),
], size=10, marker="—", bullet_color=AMBER)

# Right: future work
add_rect(s, Inches(6.75), Inches(2.05), Inches(6.0), Inches(4.95), WHITE,
         line=RGBColor(0xE3, 0xE8, 0xEF))
add_rect(s, Inches(6.75), Inches(2.05), Inches(6.0), Inches(0.5), NAVY)
add_text(s, Inches(7.05), Inches(2.05), Inches(5.4), Inches(0.5),
         "Future work — prioritized", size=14, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)

add_bullets(s, Inches(7.05), Inches(2.7), Inches(5.6), Inches(4.2), [
    ("1. Workers off Cloud Run services",
     "Move Celery to GKE with HPA on Redis depth, or convert to Cloud Run Jobs — for real scale-out"),
    ("2. Vertex AI quota lift + batched extraction",
     "Negotiate quota; explore Gemini batch API for non-real-time loads"),
    ("3. OCR for scanned PDFs",
     "Document AI integration; biggest functional gap vs the brief"),
    ("4. Cross-document consistency check",
     "Paystub income vs bank deposits within ±20% — fraud signal"),
    ("5. Phase 4 — Trusted RAG chatbot",
     "Private LLM over redacted documents, no PII exits the perimeter"),
    ("6. Multi-tenant auth + RBAC",
     "Customer / Business / Auditor roles via IAP or Auth0"),
], size=10, marker="→")

slide_footer(s, 12)


# ============================================================================
# 13. CLOSING
# ============================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, Inches(0.7), Inches(3.0), Inches(1.6), Inches(0.12), ACCENT)
add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(0.6),
         "SENTINEL", size=20, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(3.2), Inches(12), Inches(1.2),
         "Live and Reviewable", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(4.4), Inches(12), Inches(0.5),
         "sentinel-frontend-1041799394320.us-central1.run.app",
         size=15, color=RGBColor(0xCF, 0xD8, 0xE3))
add_text(s, Inches(0.7), Inches(4.85), Inches(12), Inches(0.5),
         "github.com/Khey17/sentinel-nextgenai-execution-layer",
         size=15, color=RGBColor(0xCF, 0xD8, 0xE3))
add_text(s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.5),
         "AI as assistant. Deterministic scorecard as decider.",
         size=20, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.5),
         "Every recommendation has a reason code attached.",
         size=18, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
         "Thank you  ·  Questions?",
         size=12, color=RGBColor(0x9F, 0xAD, 0xC0))


# ── Save ─────────────────────────────────────────────────────────────────────
out = "/Users/karth/Downloads/sentinel-nextgenai-execution-layer/Sentinel_Deliverable.pptx"
prs.save(out)
print("Wrote:", out)
print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst))
